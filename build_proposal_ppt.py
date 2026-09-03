"""
Build Population_DTI_INR_current_research_proposal.pptx

14 slides: split crowded finding-1 into quant + maps;
merge training/adaptation and protocol/comparison.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(r"e:\BaiduNetdiskDownload\Population-DTI-INR")
ASSETS = ROOT / "ppt_figures" / "assets"
OUT = ROOT / "Population_DTI_INR_current_research_proposal.pptx"
OUT_ALT = ROOT / "Population_DTI_INR_current_research_proposal_v5.pptx"

NAVY = RGBColor(0x16, 0x3A, 0x5F)
BLUE = RGBColor(0x4C, 0x9A, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xE8, 0xF1, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF7, 0xFA, 0xFC)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_CN = "微软雅黑"
FONT_EN = "Arial"


def _set_run_font(run, *, size_pt: float, bold: bool = False, color=NAVY, name: str | None = None, east_asia: str | None = None):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name or FONT_EN
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", east_asia or FONT_CN)


def add_title(slide, text: str, *, left=0.5, top=0.22, width=12.3, height=0.7, size=28):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run_font(run, size_pt=size, bold=True, color=NAVY)
    return box


def add_text(slide, text: str, *, left, top, width, height, size=18, bold=False, color=GRAY, align=PP_ALIGN.LEFT, font=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, size_pt=size, bold=bold, color=color, name=font or FONT_EN)
    return box


def add_subtitle(slide, text: str, *, left=0.5, top=0.85, width=12.3, height=0.45, size=14):
    return add_text(slide, text, left=left, top=top, width=width, height=height, size=size, color=GRAY, font=FONT_CN)


def add_footer(slide, page: int, total: int = 15):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.28), SLIDE_W, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    add_text(slide, f"Population-DTI-INR  ·  Research Proposal  ·  {page}/{total}", left=0.5, top=7.32, width=10, height=0.25, size=10, color=GRAY)


def add_picture_fit(slide, path: Path, left, top, width, height):
    path = Path(path)
    if not path.exists():
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        sh.fill.solid()
        sh.fill.fore_color.rgb = LIGHT
        sh.line.color.rgb = NAVY
        return sh
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def add_card(slide, left, top, width, height, fill=SOFT, line_color=NAVY):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line_color
    sh.line.width = Pt(1.25)
    return sh


def add_takeaway(slide, text: str, *, left=0.5, top=6.35, width=12.3, height=0.7):
    add_card(slide, left, top, width, height, fill=LIGHT)
    add_text(slide, text, left=left + 0.2, top=top + 0.12, width=width - 0.4, height=height - 0.15, size=14, bold=True, color=NAVY, font=FONT_CN)


def add_table(slide, rows, cols, left, top, width, height, data, col_widths=None):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if r == 0 or c > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(data[r][c])
            is_header = r == 0
            _set_run_font(run, size_pt=15 if is_header else 14, bold=is_header or c == 0, color=WHITE if is_header else NAVY)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = NAVY if is_header else (SOFT if r % 2 == 0 else WHITE)
    return table_shape


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def slide_01_title(prs):
    s = blank_slide(prs)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(s, "研究生阶段研究方案汇报", left=0.7, top=1.15, width=7, height=0.4, size=14, bold=True, color=BLUE, font=FONT_CN)
    add_text(s, "基于物理自监督隐式神经表示的\n多被试DTI参数估计研究", left=0.7, top=1.65, width=7.6, height=1.5, size=30, bold=True, color=NAVY, font=FONT_CN)
    add_text(s, "Population-DTI-INR\n面向稀疏 dMRI 采样的跨被试 DTI 参数估计", left=0.7, top=3.35, width=7.6, height=0.9, size=17, color=GRAY, font=FONT_CN)
    add_text(s, "核心问题：稀疏采样下，仅拟合信号不足以保证\nDTI 参数（如 FA）可辨识与稳定估计。", left=0.7, top=4.5, width=7.6, height=0.9, size=15, color=NAVY, font=FONT_CN)
    add_picture_fit(s, ASSETS / "title_pipeline.png", 8.6, 1.4, 4.0, 4.5)
    add_text(s, "右侧：从 dMRI 到微结构参数的总体研究链路", left=8.6, top=6.05, width=4.2, height=0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER, font=FONT_CN)
    add_footer(s, 1)


def slide_02_background(prs):
    s = blank_slide(prs)
    add_title(s, "Sparse dMRI acquisition limits reliable DTI estimation", size=26)
    add_subtitle(s, "背景：临床/加速采集常用稀疏 dMRI，测量变少后参数估计更不稳定。")
    add_picture_fit(s, ASSETS / "sparse_sampling.png", 0.4, 1.4, 5.8, 3.0)
    add_picture_fit(s, ASSETS / "problem_chain.png", 7.5, 1.35, 5.3, 4.6)
    add_card(s, 0.4, 4.55, 5.8, 2.35, fill=SOFT)
    add_text(
        s,
        "在讲什么？\n• 全采样：方向多，张量估计更稳\n• 稀疏采样：方向少，同一组信号可能对应多种 D\n• 结果：信号不确定 → 张量歧义 → FA/MD 不稳\n\n因此：稀疏条件下“能拟合信号”≠“参数正确”",
        left=0.55, top=4.7, width=5.5, height=2.1, size=13, color=NAVY, font=FONT_CN,
    )
    add_footer(s, 2)


def slide_03_methods(prs):
    s = blank_slide(prs)
    add_title(s, "Current methods still face limitations under sparse sampling", size=26)
    add_subtitle(s, "现有路线各有优势，但在稀疏采样下都难同时兼顾可解释、无标签与参数可辨识。")
    cards = [
        ("WLS", "传统物理拟合", "✓ 物理可解释\n✓ 无需深度学习", "× 稀疏时不稳定\n× 对欠采样敏感"),
        ("Deep Learning", "监督学习映射", "✓ 拟合能力强\n✓ 可端到端预测", "× 依赖 GT 参数\n× 标注成本高"),
        ("INR / QINR", "物理自监督", "✓ 无需 GT 张量\n✓ 用信号监督", "× 参数可辨识不足\n× 易出现张量漂移"),
    ]
    for i, (title, sub, good, bad) in enumerate(cards):
        left = 0.45 + i * 4.2
        add_card(s, left, 1.4, 4.0, 4.2, fill=SOFT)
        add_text(s, title, left=left + 0.15, top=1.55, width=3.7, height=0.4, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, sub, left=left + 0.15, top=2.05, width=3.7, height=0.35, size=13, color=BLUE, align=PP_ALIGN.CENTER, font=FONT_CN)
        add_text(s, good, left=left + 0.3, top=2.6, width=3.5, height=1.1, size=14, color=NAVY, font=FONT_CN)
        add_text(s, bad, left=left + 0.3, top=3.9, width=3.5, height=1.1, size=14, color=ACCENT, font=FONT_CN)
    add_takeaway(s, "本页结论：Signal fitting ≠ Parameter accuracy —— 这是后续实验要验证的核心动机。")
    add_footer(s, 3)


def slide_04_single_qinr(prs):
    s = blank_slide(prs)
    add_title(s, "Physics-driven INR enables self-supervised DTI estimation", size=26)
    add_subtitle(s, "Single-QINR：用坐标网络预测 S0 与扩散张量 D，再经 DTI 前向模型拟合观测信号。")
    add_picture_fit(s, ASSETS / "single_qinr_pipeline.png", 0.4, 1.35, 5.2, 5.5)
    add_card(s, 5.9, 1.4, 6.9, 5.4, fill=SOFT)
    add_text(
        s,
        "方法读法（按箭头往下看）\n\n1. 输入体素坐标 (x,y,z)\n2. Fourier 特征编码提高高频表达\n3. MLP 输出 S0 与张量 D\n4. 用 DTI 物理前向模型生成预测信号\n5. 与观测 DWI 做 MSE —— 物理自监督\n\n关键点\n• 训练不需要 FA/MD 真值标签\n• 监督信号来自“预测信号 vs 观测信号”\n• 因此网络可能把信号拟合得很好，\n  但张量参数仍可不唯一（下一页验证）",
        left=6.15, top=1.6, width=6.5, height=5.0, size=14, color=NAVY, font=FONT_CN,
    )
    add_footer(s, 4)


def slide_05_setup(prs):
    s = blank_slide(prs)
    add_title(s, "Single-QINR preliminary evaluation", size=28)
    add_subtitle(s, "先说明实验设置，保证后续结论可复核：在 50% 稀疏采样下评估信号拟合与 FA 一致性。")
    data = [
        ["项目", "设置"],
        ["Dataset", "HCP-YA"],
        ["Subjects", "3 名 focus 被试（106319 / 120717 / 121618）"],
        ["Sampling", "50% DWI volumes（稀疏采样）"],
        ["Shell", "b0 + b=1000"],
        ["Model", "Single-QINR（Independent INR）"],
        ["Objective", "仅信号重建（物理自监督，无 FA loss）"],
        ["Reference", "WLS-DTI（作为参数对照参考）"],
    ]
    add_table(s, 8, 2, 1.6, 1.45, 10.0, 4.6, data, col_widths=[2.6, 7.4])
    add_takeaway(s, "本页作用：告诉导师“证据从哪来”。下一页展示：信号拟合好，但 FA 仍不稳定。")
    add_footer(s, 5)


def slide_06_finding1_quant(prs):
    s = blank_slide(prs)
    add_title(s, "Key finding 1a: Signal reconstruction is strong, FA agreement is not", size=24)
    add_subtitle(s, "定量证据：信号 RelMSE 低；参数中 FA 与 WLS 一致性最弱（MD/AD/RD 一并计算）。")
    add_picture_fit(s, ASSETS / "signal_vs_parameter_gap.png", 0.35, 1.25, 7.3, 4.0)

    add_card(s, 7.85, 1.25, 5.05, 2.35, fill=SOFT)
    add_text(
        s,
        "怎么读左图？\n"
        "• 蓝柱 RelMSE≈0.07：信号拟合好\n"
        "• 红柱 FA r≈0.33–0.48：远低于 0.85\n"
        "→ 信号好 ≠ 参数好",
        left=8.0, top=1.4, width=4.75, height=2.05, size=13, color=NAVY, font=FONT_CN,
    )

    # FA/MD/AD/RD Pearson r table (3 subjects, Single-QINR 50%)
    data = [
        ["参数", "S1 r", "S2 r", "S3 r", "均值"],
        ["FA", "0.48", "0.33", "0.45", "0.42"],
        ["MD", "0.74", "0.70", "0.65", "0.70"],
        ["AD", "0.70", "0.68", "0.59", "0.66"],
        ["RD", "0.74", "0.70", "0.65", "0.70"],
    ]
    add_table(s, 5, 5, 7.85, 3.75, 5.05, 2.35, data, col_widths=[1.0, 1.0, 1.0, 1.0, 1.05])
    add_text(
        s,
        "上表：vs WLS 的 Pearson r（50% Single-QINR）。FA 最弱，故作主展示。",
        left=7.85, top=6.15, width=5.05, height=0.35, size=11, color=GRAY, font=FONT_CN,
    )
    add_takeaway(s, "结论：四类参数都算了；FA 对白质各向异性最敏感，也最能体现可辨识性问题。")
    add_footer(s, 6)


def slide_07_finding1_maps(prs):
    s = blank_slide(prs)
    add_title(s, "Key finding 1b: FA maps reveal spatial parameter error", size=26)
    add_subtitle(s, "空间证据：相对 WLS 参考，Single-QINR 的 FA 细节丢失，误差集中在白质。")
    add_picture_fit(s, ASSETS / "fa_maps_comparison.png", 0.4, 1.3, 8.3, 4.55)
    add_card(s, 8.9, 1.35, 4.0, 4.5, fill=SOFT)
    add_text(
        s,
        "三列怎么看？\n\n1) Reference FA (WLS)\n   白质结构清晰\n\n2) Single-QINR FA\n   细节模糊、对比度下降\n\n3) Absolute Error\n   高误差集中在白质束\n\n被试 106319\n• FA r ≈ 0.48\n• FA MAE ≈ 0.144\n\n与上一页数字一致：\n参数层面确实不稳。",
        left=9.1, top=1.5, width=3.65, height=4.25, size=13, color=NAVY, font=FONT_CN,
    )
    add_takeaway(s, "定量（含 FA/MD/AD/RD）+ 空间 FA 图共同支持：稀疏采样下存在参数可辨识性问题。")
    add_footer(s, 7)


def slide_08_finding2(prs):
    s = blank_slide(prs)
    add_title(s, "Key finding 2: Local regularization cannot remove ambiguity", size=24)
    add_subtitle(s, "Phase 9：加大张量正则 λ，信号可改善，但 FA 误差几乎不变。")
    add_picture_fit(s, ASSETS / "phase9_regularization.png", 0.3, 1.3, 7.9, 4.55)
    add_card(s, 8.4, 1.35, 4.5, 4.5, fill=SOFT)
    add_text(
        s,
        "Phase 9 在做什么？\n信号适应时加入不同强度 λ_dis，\n看 FA 是否稳住。\n\n结果\n• λ ∈ {0, 0.01, 0.1}\n• FA MAE ≈ 0.175（平坦）\n• 平均 ΔPSNR ≈ +11 dB\n• 信号提升，FA 不随之恢复\n\n说明\n仅靠局部正则不够，\n需要跨被试先验约束。",
        left=8.6, top=1.5, width=4.15, height=4.2, size=13, color=NAVY, font=FONT_CN,
    )
    add_takeaway(s, "本页结论：参数歧义不能只靠局部正则解决 → 需要更强先验（见文献动机）。")
    add_footer(s, 8)


def slide_09_literature(prs):
    """Bridge: SM-INR + CINA motivate Population-DTI-INR."""
    s = blank_slide(prs)
    add_title(s, "Literature motivation: why move to a population prior?", size=24)
    add_subtitle(s, "两篇相关工作分别提供“单被试 INR 可行”与“跨被试共享+latent 范式”。")

    # Left: SM-INR
    add_card(s, 0.35, 1.35, 6.2, 3.55, fill=SOFT)
    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(1.35), Inches(6.2), Inches(0.55))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = NAVY
    hdr.line.fill.background()
    add_text(
        s, "SM-INR（Communications Biology, 2026）",
        left=0.45, top=1.45, width=6.0, height=0.4, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_CN,
    )
    add_text(
        s,
        "Implicit neural representations for accurate\nestimation of the Standard Model of white matter\n\n"
        "要点\n"
        "• 用 INR 自监督拟合 dMRI 信号\n"
        "• 估计白质 Standard Model 微结构参数\n"
        "• 空间连续编码提供单被试内正则\n"
        "• 指出高维微结构估计存在 degeneracy\n\n"
        "对我们：支持“物理自监督 INR 估微结构”",
        left=0.55, top=2.05, width=5.85, height=2.7, size=12, color=NAVY, font=FONT_CN,
    )

    # Right: CINA
    add_card(s, 6.8, 1.35, 6.15, 3.55, fill=LIGHT)
    hdr2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.35), Inches(6.15), Inches(0.55))
    hdr2.fill.solid()
    hdr2.fill.fore_color.rgb = BLUE
    hdr2.line.fill.background()
    add_text(
        s, "CINA（MICCAI / arXiv 2024）",
        left=6.9, top=1.45, width=5.95, height=0.4, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_CN,
    )
    add_text(
        s,
        "Conditional Implicit Neural Atlas for\nSpatio-Temporal Representation of Fetal Brains\n\n"
        "要点\n"
        "• 共享网络 θ 学群体共性（atlas）\n"
        "• 每人 latent z 编码个体差异\n"
        "• 新被试：冻 θ，只优化 z（test-time）\n"
        "• auto-decoder 式跨被试 INR\n\n"
        "对我们：提供 Population 框架模板",
        left=7.0, top=2.05, width=5.8, height=2.7, size=12, color=NAVY, font=FONT_CN,
    )

    # Bottom logic chain
    add_card(s, 0.35, 5.1, 12.6, 1.8, fill=SOFT)
    add_text(
        s,
        "合在一起的逻辑链（讲解用）\n"
        "SM-INR：单被试 INR 可估微结构  →  我们的实验：稀疏下信号好、FA 仍差，局部约束无效\n"
        "→ 需要更强先验  →  CINA：跨被试共享 θ + 个体 z  →  Population-DTI-INR\n"
        "= CINA 式群体框架  +  SM-INR 式物理自监督张量/微结构估计",
        left=0.55, top=5.25, width=12.2, height=1.5, size=13, color=NAVY, font=FONT_CN,
    )
    add_footer(s, 9)


def slide_10_hypothesis(prs):
    s = blank_slide(prs)
    add_title(s, "Can population prior improve DTI parameter identifiability?", size=26)
    add_subtitle(s, "研究假设：共享群体表示 + 被试潜变量，能否约束稀疏采样下的参数可辨识性？")
    add_picture_fit(s, ASSETS / "single_vs_population.png", 0.4, 1.35, 8.4, 4.5)
    add_card(s, 9.0, 1.4, 3.9, 4.4, fill=SOFT)
    add_text(
        s,
        "左右对比\n\n左：Single\n每被试单独训 INR\n缺少跨被试约束\n\n右：Population\n共享主干 θ + latent z\n\n假设\n群体先验可缩小可行\n张量解空间，提升稳定性。",
        left=9.2, top=1.55, width=3.55, height=4.15, size=13, color=NAVY, font=FONT_CN,
    )
    add_takeaway(s, "不是换热闹方法：实验发现问题 + SM-INR/CINA 提供路径 → Population-DTI-INR。")
    add_footer(s, 10)


def slide_11_architecture(prs):
    s = blank_slide(prs)
    add_title(s, "Population-DTI-INR framework", size=28)
    add_subtitle(s, "共享参数 θ 编码群体共性；被试 latent z_s 编码个体差异（呼应 CINA）。")
    add_picture_fit(s, ASSETS / "population_architecture.png", 0.35, 1.3, 8.5, 4.7)
    add_card(s, 9.05, 1.35, 3.85, 4.6, fill=SOFT)
    add_text(
        s,
        "模块说明\n\n• Subject ID → z_s\n• Coordinate → Fourier\n• Shared MLP (θ)\n• S0 head / D head\n• DTI Forward → 预测信号\n\n深蓝 = 共享 θ\n浅蓝 = 被试 z",
        left=9.25, top=1.55, width=3.5, height=4.3, size=14, color=NAVY, font=FONT_CN,
    )
    add_takeaway(s, "训练仍走物理自监督（呼应 SM-INR），不直接回归 FA。")
    add_footer(s, 11)


def slide_12_train_adapt(prs):
    s = blank_slide(prs)
    add_title(s, "Training and generalization to unseen subjects", size=26)
    add_subtitle(s, "训练只用信号 MSE；未见被试用 zero-shot / latent adaptation（呼应 CINA test-time）。")
    add_picture_fit(s, ASSETS / "training_strategy.png", 0.35, 1.3, 6.2, 3.5)
    add_picture_fit(s, ASSETS / "subject_adaptation.png", 6.7, 1.3, 6.2, 3.5)
    add_card(s, 0.35, 4.95, 6.2, 1.95, fill=SOFT)
    add_text(
        s,
        "训练（左）\n• Loss = MSE(S_pred, S_obs)\n• 优化 θ + z_s\n• 不用 GT tensor / FA loss\n→ 保持物理自监督设定",
        left=0.5, top=5.1, width=5.95, height=1.7, size=13, color=NAVY, font=FONT_CN,
    )
    add_card(s, 6.7, 4.95, 6.2, 1.95, fill=LIGHT)
    add_text(
        s,
        "未见被试（右）\n• Zero-shot：θ 冻结，z_new=0\n  → 检验群体先验本身\n• Adaptation：θ 冻结，只优化 z\n  → 检验个体适配增益",
        left=6.85, top=5.1, width=5.95, height=1.7, size=13, color=NAVY, font=FONT_CN,
    )
    add_footer(s, 12)


def slide_13_protocol_compare(prs):
    s = blank_slide(prs)
    add_title(s, "Evaluation protocol and comparison methods", size=26)
    add_subtitle(s, "同一稀疏设定下，用三种方法回答先验与适配各自贡献多少。")
    data_left = [
        ["项目", "设置"],
        ["Dataset", "HCP-YA"],
        ["Split", "Train 3 / Val 1 / Test 1"],
        ["Sampling", "50% DWI"],
        ["Shell", "b0 + b1000"],
        ["Metrics", "Signal + Parameter + Stability"],
    ]
    add_table(s, 6, 2, 0.4, 1.4, 5.9, 3.6, data_left, col_widths=[1.8, 4.1])
    data_right = [
        ["方法", "回答的问题"],
        ["Single-QINR", "单被试自监督上限？"],
        ["Pop. zero-shot", "群体先验本身有用吗？"],
        ["Pop. adaptation", "个体适配能否稳住参数？"],
    ]
    add_table(s, 4, 2, 6.6, 1.4, 6.3, 2.8, data_right, col_widths=[2.6, 3.7])
    add_card(s, 6.6, 4.4, 6.3, 2.45, fill=SOFT)
    add_text(
        s,
        "汇报时怎么讲？\n先固定数据与采样协议，再比较三种方法。\n重点不是只看 PSNR，而是 FA / 漂移是否同步改善。\nPopulation 完整数字表见下一页“待补”栏。",
        left=6.8, top=4.55, width=6.0, height=2.2, size=13, color=NAVY, font=FONT_CN,
    )
    add_card(s, 0.4, 5.2, 5.9, 1.65, fill=LIGHT)
    add_text(
        s,
        "三类指标：\nSignal（RMSE/PSNR/RelMSE）\nParameter（FA/MD/AD/RD）\nStability（tensor drift / ΔFA）",
        left=0.55, top=5.35, width=5.6, height=1.4, size=13, color=NAVY, font=FONT_CN,
    )
    add_footer(s, 13)


def slide_14_metrics(prs):
    s = blank_slide(prs)
    add_title(s, "Preliminary numbers vs pending Population evaluation", size=24)
    add_subtitle(s, "已有数字支撑“问题成立”；Population 完整对比仍在后续验证。")
    add_card(s, 0.35, 1.35, 6.25, 5.5, fill=SOFT)
    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(1.35), Inches(6.25), Inches(0.65))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = NAVY
    hdr.line.fill.background()
    add_text(s, "已有（支撑研究问题）", left=0.45, top=1.48, width=6.05, height=0.4, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_CN)
    add_text(
        s,
        "Single-QINR · 50% · 3 被试\n• RelMSE ≈ 0.070–0.076  → 信号好\n• FA r ≈ 0.33–0.48       → 低于 0.85\n• FA MAE ≈ 0.14–0.15\n\nPhase 9 · λ 消融 · 4 被试\n• FA MAE ≈ 0.175（几乎不随 λ 变）\n• holdout ΔPSNR ≈ +11 dB\n• 信号变好，FA 未必变好\n\n这些数字回答：问题真实存在。",
        left=0.55, top=2.15, width=5.9, height=4.5, size=13, color=NAVY, font=FONT_CN,
    )
    add_card(s, 6.85, 1.35, 6.05, 5.5, fill=LIGHT)
    hdr2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.85), Inches(1.35), Inches(6.05), Inches(0.65))
    hdr2.fill.solid()
    hdr2.fill.fore_color.rgb = BLUE
    hdr2.line.fill.background()
    add_text(s, "待补（Population 完整验证）", left=6.95, top=1.48, width=5.85, height=0.4, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_CN)
    add_text(
        s,
        "下一步统一对比表：\n• Single-QINR\n• Population zero-shot\n• Population adaptation\n\n对齐报告：\n• Signal：RMSE / PSNR / RelMSE\n• Parameter：FA / MD / AD / RD\n• Stability：drift、ΔFA\n\n完整可汇报对比尚在验证，\n避免过早下成功结论。",
        left=7.05, top=2.15, width=5.7, height=4.5, size=13, color=NAVY, font=FONT_CN,
    )
    add_footer(s, 14)


def slide_15_roadmap(prs):
    s = blank_slide(prs)
    add_title(s, "Research roadmap", size=28)
    add_subtitle(s, "实验发现问题 → 文献提供路径 → Population-DTI-INR 验证。")
    add_picture_fit(s, ASSETS / "research_roadmap.png", 0.5, 1.35, 12.3, 2.5)
    add_card(s, 0.5, 4.05, 12.3, 2.85, fill=SOFT)
    add_text(
        s,
        "向导师汇报时的一句话逻辑\n\n"
        "1）Sparse 下 Single-QINR：信号好、FA 不稳（定量 + FA 图）\n"
        "2）局部正则无法消除歧义（Phase 9）\n"
        "3）文献：SM-INR 支持物理自监督 INR；CINA 提供共享θ+z范式\n"
        "4）因此提出 Population-DTI-INR，用 zero-shot / adaptation 验证 FA 可辨识性\n"
        "5）后续可扩展到 DKI 等更高阶模型\n\n"
        "总结：Population prior 有望改善稀疏 dMRI 下的 DTI 参数可辨识性。",
        left=0.75, top=4.2, width=11.9, height=2.6, size=13, color=NAVY, font=FONT_CN,
    )
    add_footer(s, 15)


FORBIDDEN = ["fail", "failed", "failure", "wrong", "abandon", "abandoned", "change direction", "原开题", "方向调整", "放弃", "错误", "失败"]


def validate_pptx(prs: Presentation) -> list[str]:
    issues = []
    texts = []
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                texts.append(t)
                low = t.lower()
                for bad in FORBIDDEN:
                    if bad.lower() in low or bad in t:
                        issues.append(f"Slide {i}: forbidden term '{bad}'")
    if len(prs.slides) > 15:
        issues.append(f"Too many slides: {len(prs.slides)}")
    joined = " ".join(texts).lower()
    for need in ["signal", "parameter", "population", "identifiab"]:
        if need not in joined:
            issues.append(f"Missing theme keyword: {need}")
    return issues


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_01_title(prs)
    slide_02_background(prs)
    slide_03_methods(prs)
    slide_04_single_qinr(prs)
    slide_05_setup(prs)
    slide_06_finding1_quant(prs)
    slide_07_finding1_maps(prs)
    slide_08_finding2(prs)
    slide_09_literature(prs)
    slide_10_hypothesis(prs)
    slide_11_architecture(prs)
    slide_12_train_adapt(prs)
    slide_13_protocol_compare(prs)
    slide_14_metrics(prs)
    slide_15_roadmap(prs)
    issues = validate_pptx(prs)
    try:
        prs.save(str(OUT))
        saved = OUT
    except PermissionError:
        prs.save(str(OUT_ALT))
        saved = OUT_ALT
        print(f"NOTE: original locked; saved as {saved.name}")
    print(f"Saved: {saved}")
    print(f"Slides: {len(prs.slides)}")
    print("VALIDATION:", "PASS" if not issues else issues)


if __name__ == "__main__":
    main()
